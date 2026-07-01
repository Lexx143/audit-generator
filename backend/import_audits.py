import os
import glob
import json
import uuid
from pathlib import Path
from pptx import Presentation
from pypdf import PdfReader
from dotenv import load_dotenv
from openai import OpenAI
import chromadb
from pydantic import BaseModel
from typing import List, Optional

load_dotenv()

class AuditCase(BaseModel):
    title: str
    category: str
    priority: str
    risk: str
    recommendation: Optional[str] = None
    vulnerability: str = ""

class AuditData(BaseModel):
    cases: List[AuditCase]
    conclusions: List[str]

chroma_client = chromadb.PersistentClient(path="db")
from chromadb.utils import embedding_functions

openai_ef = embedding_functions.OpenAIEmbeddingFunction(
    api_key=os.environ.get("OPENAI_API_KEY"),
    model_name="text-embedding-3-small"
)
cases_collection = chroma_client.get_or_create_collection(name="audit_cases", embedding_function=openai_ef)
conclusions_collection = chroma_client.get_or_create_collection(name="audit_conclusions", embedding_function=openai_ef)

client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

HINTS_FILE = "db/hints.json"

def get_hints():
    if os.path.exists(HINTS_FILE):
        with open(HINTS_FILE, "r") as f:
            return json.load(f)
    return []

def save_hints(hints: list):
    with open(HINTS_FILE, "w") as f:
        json.dump(hints, f, ensure_ascii=False, indent=2)

def extract_text_from_pptx(filepath):
    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_pdf(filepath):
    try:
        reader = PdfReader(filepath)
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def parse_and_save(text, filename):
    if not text.strip():
        return
        
    print(f"Processing {filename}...")
    prompt = f"""
    Извлеки все ИТ-кейсы, проблемы, уязвимости, риски и рекомендации из этого текста (это ИТ-аудит).
    Если находишь конкретную проблему (кейс), заполняй все поля (category, priority, risk, recommendation).
    priority используй 'high', 'medium' или 'low'.
    Верни список кейсов и общих выводов.
    
    ТЕКСТ АУДИТА:
    {text[:30000]}
    """
    
    try:
        response = client.beta.chat.completions.parse(
            model='gpt-5.5',
            messages=[
                {"role": "system", "content": "Ты эксперт по ИТ-аудитам. Извлекай данные в структурированном виде."},
                {"role": "user", "content": prompt}
            ],
            response_format=AuditData,
        )
        data_obj = response.choices[0].message.parsed
        data = json.loads(data_obj.model_dump_json())
        
        current_hints = get_hints()
        new_hints_added = False
        
        for case in data.get("cases", []):
            case_text = f"Тема: {case['title']}. Категория: {case['category']}. Риск: {case['risk']}. Рекомендация: {case.get('recommendation', '')}"
            cases_collection.upsert(
                documents=[case_text],
                ids=[str(uuid.uuid4())]
            )
            if case['title'] not in current_hints:
                current_hints.append(case['title'])
                new_hints_added = True
                
        for conc in data.get("conclusions", []):
            conclusions_collection.upsert(
                documents=[conc],
                ids=[str(uuid.uuid4())]
            )
            
        if new_hints_added:
            save_hints(current_hints)
            
        print(f"Saved {len(data.get('cases', []))} cases from {filename}.")
        
    except Exception as e:
        print(f"Failed to parse {filename} with OpenAI: {e}")

def main():
    base_dir = os.path.expanduser("~/Documents/Документы_Ноутбук_Александр/Астана")
    files = list(Path(base_dir).rglob("*аудит*.*"))
    files = [f for f in files if f.suffix.lower() in ['.pptx', '.pdf'] and not f.name.startswith("~$")]
    
    processed_names = set()
    
    for f in files:
        if f.stem in processed_names:
            continue
            
        processed_names.add(f.stem)
        
        if f.suffix.lower() == '.pptx':
            text = extract_text_from_pptx(f)
        else:
            text = extract_text_from_pdf(f)
            
        parse_and_save(text, f.name)
        
    print("ALL DONE!")

if __name__ == "__main__":
    main()
