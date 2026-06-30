import os
import json
import base64
import time
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from typing import List, Optional
from google import genai
from google.genai import types
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy

load_dotenv()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

os.makedirs("static/images", exist_ok=True)
app.mount("/static", StaticFiles(directory="static"), name="static")

try:
    client = genai.Client(
        vertexai=True, 
        project=os.environ.get("PROJECT_ID", "gemini-bot-460610"), 
        location=os.environ.get("LOCATION", "us-central1")
    )
except Exception as e:
    print(f"Warning: Failed to initialize Vertex AI client: {e}")
    client = None

# ---- Pydantic Models ----
class Case(BaseModel):
    title: str
    vulnerability: str
    risk: str
    recommendation: Optional[str]
    priority: str 
    category: str 
    image_b64: Optional[str] = None
    image_prompt: Optional[str] = None

class AuditData(BaseModel):
    client_name: str
    review: str 
    cases: List[Case] 
    conclusions: List[str] 

class ParseRequest(BaseModel):
    general_data: str
    vulnerabilities: str
    conclusions: str
    audit_type: str 

class GenerateImageRequest(BaseModel):
    prompt: str

class GeneratePptxRequest(BaseModel):
    data: AuditData

# ---- API Endpoints ----
@app.post("/api/parse")
async def parse_audit(req: ParseRequest):
    if not client:
        raise HTTPException(status_code=500, detail="Vertex AI client not initialized")
    
    prompt = f"""
    Проанализируй предоставленные данные и сформируй структуру для ИТ-аудита.
    Тип аудита: {req.audit_type}. Если 'full', обязательно добавь рекомендации (поле recommendation) для каждого кейса. Если 'express', оставь поле рекомендаций пустым.
    
    ОБЩИЕ ДАННЫЕ О КЛИЕНТЕ:
    {req.general_data}
    
    ВЫЯВЛЕННЫЕ УЯЗВИМОСТИ:
    {req.vulnerabilities}
    
    ВЫВОДЫ/ЗАКЛЮЧЕНИЕ:
    {req.conclusions}
    
    ИНСТРУКЦИИ ДЛЯ СТРУКТУРЫ (AuditData):
    1. client_name: Короткое название клиента.
    2. review: Обобщающий текст-обзор на 2-3 предложения.
    3. cases: Выдели ровно 5 основных уязвимостей. Для каждой заполни title, vulnerability, risk, priority ("ПЕРВЫЙ ПРИОРИТЕТ", "ВТОРОЙ ПРИОРИТЕТ" или "ТРЕТИЙ ПРИОРИТЕТ"), category ("I. Серверная инфраструктура" или "II. Сеть и ИТ-поддержка"). Также придумай image_prompt - промпт на английском (до 15 слов) для генерации ИТ-векторной картинки без текста.
    4. conclusions: Массив из 3-5 строк для итогового слайда с предложениями.
    """
    
    try:
        response = client.models.generate_content(
            model='gemini-2.5-pro',
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=AuditData,
                temperature=0.2,
            ),
        )
        data = json.loads(response.text)
        return data
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate_image")
async def generate_image(req: GenerateImageRequest):
    if not client:
        raise HTTPException(status_code=500, detail="Vertex AI client not initialized")
        
    full_prompt = f"Minimalist abstract 3D icon or simple vector of {req.prompt}. Pure white background, absolutely no text, no letters, clean corporate IT infographic style, geometric."
    try:
        response = client.models.generate_images(
            model='imagen-3.0-generate-001',
            prompt=full_prompt,
            config=types.GenerateImagesConfig(
                number_of_images=1,
                aspect_ratio="1:1"
            )
        )
        img_bytes = response.generated_images[0].image.image_bytes
        img_b64 = base64.b64encode(img_bytes).decode('utf-8')
        return {"image_b64": f"data:image/jpeg;base64,{img_b64}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def perfect_replace(shape, new_text):
    if not hasattr(shape, "text_frame"): return
    tf = shape.text_frame
    if not tf.paragraphs: return
    
    first_run = None
    first_p_idx = 0
    for i, p in enumerate(tf.paragraphs):
        if p.runs:
            first_run = p.runs[0]
            first_p_idx = i
            break
            
    if not first_run:
        shape.text = new_text
        return
        
    first_run.text = new_text
    
    for i in range(1, len(tf.paragraphs[first_p_idx].runs)):
        tf.paragraphs[first_p_idx].runs[i].text = ""
        
    for i in range(len(tf.paragraphs)):
        if i == first_p_idx: continue
        for r in tf.paragraphs[i].runs:
            r.text = ""

def replace_table_cell(table, r, c, new_text):
    if not table: return
    cell = table.cell(r, c)
    perfect_replace(cell, new_text)

def get_table(slide):
    for s in slide.shapes:
        if s.has_table: return s.table
    return None

def get_title(slide, fallback_idx=1):
    for s in slide.shapes:
        if hasattr(s, 'text') and s.text and 'Кейс' in s.text: return s
    if len(slide.shapes) > fallback_idx: return slide.shapes[fallback_idx]
    return None

def get_priority(slide):
    for s in slide.shapes:
        if hasattr(s, 'text') and s.text and 'ПРИОРИТЕТ' in s.text: return s
    return None

@app.post("/api/generate_pptx")
async def generate_pptx(req: GeneratePptxRequest):
    data = req.data
    template_path = os.path.join("assets", "template.pptx")
    output_filename = f"report_{int(time.time())}.pptx"
    output_path = os.path.join("static", output_filename)
    
    if not os.path.exists(template_path):
        raise HTTPException(status_code=500, detail="Template not found")
        
    prs = Presentation(template_path)
    
    temp_images = []
    for i, case in enumerate(data.cases):
        if case.image_b64 and case.image_b64.startswith("data:image"):
            img_data = base64.b64decode(case.image_b64.split(",")[1])
            img_path = f"/tmp/case_img_{i}_{int(time.time())}.jpg"
            with open(img_path, "wb") as f:
                f.write(img_data)
            temp_images.append((i, img_path))
            
    for s in prs.slides[0].shapes:
        if hasattr(s, 'text') and s.text:
            if 'Клиент' in s.text or 'Курылыс' in s.text: 
                perfect_replace(s, f"Клиент: {data.client_name}")
                
    for s in prs.slides[2].shapes:
        if hasattr(s, 'text') and s.text and 'аудита' in s.text:
            perfect_replace(s, data.review)
            
    prio_counts = {"ПЕРВЫЙ ПРИОРИТЕТ": 0, "ВТОРОЙ ПРИОРИТЕТ": 0, "ТРЕТИЙ ПРИОРИТЕТ": 0}
    cat_counts = {
        "I. Серверная инфраструктура": {"ПЕРВЫЙ ПРИОРИТЕТ": 0, "ВТОРОЙ ПРИОРИТЕТ": 0, "ТРЕТИЙ ПРИОРИТЕТ": 0},
        "II. Сеть и ИТ-поддержка": {"ПЕРВЫЙ ПРИОРИТЕТ": 0, "ВТОРОЙ ПРИОРИТЕТ": 0, "ТРЕТИЙ ПРИОРИТЕТ": 0}
    }
    
    for case in data.cases:
        prio_counts[case.priority] += 1
        if case.category in cat_counts:
            cat_counts[case.category][case.priority] += 1
            
    t3 = get_table(prs.slides[3])
    if t3:
        replace_table_cell(t3, 0, 0, str(prio_counts["ПЕРВЫЙ ПРИОРИТЕТ"]))
        replace_table_cell(t3, 0, 1, str(prio_counts["ВТОРОЙ ПРИОРИТЕТ"]))
        replace_table_cell(t3, 0, 2, str(prio_counts["ТРЕТИЙ ПРИОРИТЕТ"]))

    t4 = get_table(prs.slides[4])
    if t4:
        cats = list(cat_counts.keys())
        for r, cat in enumerate(cats):
            replace_table_cell(t4, r, 0, cat)
            replace_table_cell(t4, r, 1, str(cat_counts[cat]["ПЕРВЫЙ ПРИОРИТЕТ"]))
            replace_table_cell(t4, r, 2, str(cat_counts[cat]["ВТОРОЙ ПРИОРИТЕТ"]))
            replace_table_cell(t4, r, 3, str(cat_counts[cat]["ТРЕТИЙ ПРИОРИТЕТ"]))
        
        for r in range(len(cats), 3):
            replace_table_cell(t4, r, 0, "")
            replace_table_cell(t4, r, 1, "")
            replace_table_cell(t4, r, 2, "")
            replace_table_cell(t4, r, 3, "")
            
    for s in prs.slides[5].shapes:
        if hasattr(s, 'text') and s.text and 'I.' in s.text:
            perfect_replace(s, "I. Серверная инфраструктура")
            
    for s in prs.slides[9].shapes:
        if hasattr(s, 'text') and s.text and 'II.' in s.text:
            perfect_replace(s, "II. Сеть и ИТ-поддержка")

    case_slides = [6, 7, 8, 10, 11]
    for i, case in enumerate(data.cases):
        if i >= len(case_slides): break
        slide = prs.slides[case_slides[i]]
        perfect_replace(get_title(slide), f"Кейс {i+1}: {case.title}")
        
        if get_priority(slide):
            perfect_replace(get_priority(slide), case.priority)
            
        t = get_table(slide)
        if t:
            replace_table_cell(t, 0, 0, f"Уязвимость: {case.vulnerability}")
            risk_text = f"Риски: {case.risk}"
            if case.recommendation:
                risk_text += f"\n\nРекомендации: {case.recommendation}"
            replace_table_cell(t, 1, 0, risk_text)
            
        img_path_tuple = next((x for x in temp_images if x[0] == i), None)
        if img_path_tuple:
            for s in slide.shapes:
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    l, t_pos, w, h = s.left, s.top, s.width, s.height
                    s._element.getparent().remove(s._element)
                    slide.shapes.add_picture(img_path_tuple[1], l, t_pos, w, h)
                    break

    conc_text = "\n".join([f"• {c}" for c in data.conclusions])
    for s in prs.slides[16].shapes:
        if hasattr(s, 'text') and s.text and 'Вариант 1' in s.text:
            perfect_replace(s, conc_text)
            
    for i in sorted([15, 14, 13, 12], reverse=True):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
        
    prs.save(output_path)
    
    for _, img_path in temp_images:
        if os.path.exists(img_path):
            os.remove(img_path)
            
    return {"url": f"http://localhost:8000/static/{output_filename}"}
