import base64
import copy
import hashlib
import io
import math
import os
import re
import time
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from schemas import AuditData

TEMPLATE_PATH = os.path.join("assets", "template.pptx")

# Геометрия слайда кейса (портрет 7559675 x 10691813 EMU):
# таблица top=2444005, блок "Ответ ИТ-руководства" top~8339406, картинка в шаблоне 3111500.
CASE_TABLE_CHARS_PER_LINE = 55   # ~14pt на ширине таблицы 6.26"
CASE_TABLE_LINE_EMU = 220_000    # высота строки текста при 14pt
CASE_TABLE_ROW_PAD_EMU = 100_000
CASE_IMAGE_MARGIN_EMU = 150_000
CASE_IMAGE_MAX_EMU = 3_111_500
CASE_IMAGE_MIN_EMU = 1_100_000

# Фирменные цвета плашек приоритета (из шаблона)
PRIORITY_COLORS = {
    "ПЕРВЫЙ ПРИОРИТЕТ": RGBColor(0xE3, 0x4A, 0x4E),
    "ВТОРОЙ ПРИОРИТЕТ": RGBColor(0xFB, 0xBA, 0x36),
    "ТРЕТИЙ ПРИОРИТЕТ": RGBColor(0x01, 0xB9, 0xD2),
}

COMPANY_ADDRESS = "Янтарная улица, 58в"

ROMAN = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X",
         "XI", "XII", "XIII", "XIV", "XV", "XVI", "XVII", "XVIII", "XIX", "XX"]


def _strip_category_number(cat: str) -> str:
    return re.sub(r"^\s*[IVXivx]+\.\s*", "", (cat or "").strip())


def perfect_replace(shape, new_text):
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return

    first_run = None
    first_p_idx = 0
    for i, p in enumerate(tf.paragraphs):
        if p.runs:
            first_run = p.runs[0]
            first_p_idx = i
            break

    if not first_run:
        shape.text = new_text
        return

    first_run.text = new_text

    for i in range(1, len(tf.paragraphs[first_p_idx].runs)):
        tf.paragraphs[first_p_idx].runs[i].text = ""

    for i in range(len(tf.paragraphs)):
        if i == first_p_idx:
            continue
        for r in tf.paragraphs[i].runs:
            r.text = ""


def replace_table_cell(table, r, c, new_text):
    if not table:
        return
    cell = table.cell(r, c)
    perfect_replace(cell, new_text)


def replace_labeled_cell(table, r, c, label, body):
    """Заполняет ячейку 'Метка: текст', сохраняя структуру шаблона:
    первый run — жирная метка, второй — обычный текст."""
    cell = table.cell(r, c)
    paragraphs = cell.text_frame.paragraphs
    first_idx = None
    for i, p in enumerate(paragraphs):
        if p.runs:
            first_idx = i
            break

    if first_idx is None or len(paragraphs[first_idx].runs) < 2:
        perfect_replace(cell, f"{label}{body}")
        return

    first_p = paragraphs[first_idx]
    first_p.runs[0].text = label
    first_p.runs[1].text = body
    for run in first_p.runs[2:]:
        run.text = ""
    for i, p in enumerate(paragraphs):
        if i == first_idx:
            continue
        for run in p.runs:
            run.text = ""


def get_table(slide):
    for s in slide.shapes:
        if s.has_table:
            return s.table
    return None


def get_table_shape(slide):
    for s in slide.shapes:
        if s.has_table:
            return s
    return None


def get_title(slide, fallback_idx=1):
    for s in slide.shapes:
        if hasattr(s, 'text') and s.text and 'Кейс' in s.text:
            return s
    if len(slide.shapes) > fallback_idx:
        return slide.shapes[fallback_idx]
    return None


def get_priority(slide):
    for s in slide.shapes:
        if hasattr(s, 'text') and s.text and 'ПРИОРИТЕТ' in s.text:
            return s
    return None


def recolor_priority_badge(slide, priority):
    """Красит плашку под текстом приоритета в цвет, соответствующий приоритету."""
    color = PRIORITY_COLORS.get(priority)
    prio_box = get_priority(slide)
    if color is None or prio_box is None:
        return
    cx = prio_box.left + prio_box.width // 2
    cy = prio_box.top + prio_box.height // 2
    for s in slide.shapes:
        if s.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if s.left <= cx <= s.left + s.width and s.top <= cy <= s.top + s.height:
            s.fill.solid()
            s.fill.fore_color.rgb = color


def _delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def clone_slide(prs, source):
    """Клонирует слайд (фигуры + картинки) в конец презентации."""
    dest = prs.slides.add_slide(source.slide_layout)
    # add_slide подтягивает плейсхолдеры макета — убираем, фигуры скопируем сами
    for shp in list(dest.shapes):
        _delete_shape(shp)

    # переносим связи на картинки и запоминаем соответствие rId
    rid_map = {}
    for rel in source.part.rels.values():
        if rel.is_external or "image" not in rel.reltype:
            continue
        rid_map[rel.rId] = dest.part.relate_to(rel.target_part, rel.reltype)

    for shp in source.shapes:
        el = copy.deepcopy(shp._element)
        for sub in el.iter():
            for key in list(sub.attrib):
                if (key.endswith('}embed') or key.endswith('}link')) and sub.attrib[key] in rid_map:
                    sub.set(key, rid_map[sub.attrib[key]])
        dest.shapes._spTree.append(el)
    return dest


def _image_md5(shape):
    try:
        return hashlib.md5(shape.image.blob).hexdigest()
    except Exception:
        return None


def remove_personal_marks(prs, auditor=None):
    """Заменяет данные аудитора из шаблона на выбранного аудитора
    (или убирает их совсем, если аудитор не указан)."""
    auditor_name = (auditor.name.strip() if auditor and auditor.name else None)
    photo_bytes = None
    if auditor and auditor.photo_b64:
        photo_bytes = base64.b64decode(auditor.photo_b64.split(",", 1)[-1])

    title_slide = prs.slides[0]

    # Фото на титуле — картинка в нижней трети слайда (лого сидит вверху)
    photo_md5 = None
    for s in list(title_slide.shapes):
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and s.top and s.top > prs.slide_height * 0.6:
            photo_md5 = _image_md5(s)
            geom = (s.left, s.top, s.width, s.height)
            _delete_shape(s)
            if photo_bytes:
                title_slide.shapes.add_picture(io.BytesIO(photo_bytes), *geom)

    # То же фото на других слайдах (в т.ч. picture placeholder на слайде "Ревью")
    if photo_md5:
        for slide in prs.slides:
            for s in list(slide.shapes):
                if s.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER):
                    if _image_md5(s) == photo_md5:
                        geom = (s.left, s.top, s.width, s.height)
                        _delete_shape(s)
                        if photo_bytes:
                            slide.shapes.add_picture(io.BytesIO(photo_bytes), *geom)

    # Титул: "Отчет сформирован: <ФИО>" и текущая дата вместо заглушек шаблона
    for s in title_slide.shapes:
        if not (hasattr(s, 'text') and s.text and 'сформирован' in s.text):
            continue
        paragraphs = s.text_frame.paragraphs
        for i, p in enumerate(paragraphs):
            new_text = None
            if 'сформирован' in p.text:
                new_text = "Отчет сформирован: " + auditor_name if auditor_name else "Отчет сформирован"
            elif 'Дата' in p.text:
                new_text = f"Дата: {date.today().strftime('%d.%m.%Y')}"
            if new_text is None:
                continue
            for ri, r in enumerate(p.runs):
                r.text = new_text if ri == 0 else ""

    # Внутренний номер: "+7 ... вн. 162" -> оставляем только основной номер
    for slide in prs.slides:
        for s in slide.shapes:
            if not (hasattr(s, 'text') and s.text and '162' in s.text and 'вн' in s.text):
                continue
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if 'вн' in r.text or '162' in r.text:
                        r.text = ''

    # Блок "Отчет составил <ФИО, должность, личный телефон>"
    for slide in prs.slides:
        for s in list(slide.shapes):
            if not (hasattr(s, 'text') and s.text and 'Отчет составил' in s.text):
                continue
            if not auditor_name:
                _delete_shape(s)
                continue
            # Первый абзац "Отчет составил" и второй с ФИО оставляем,
            # должность и личный телефон из шаблона затираем
            paragraphs = s.text_frame.paragraphs
            for i, p in enumerate(paragraphs):
                for ri, r in enumerate(p.runs):
                    if i == 0:
                        continue
                    r.text = auditor_name if (i == 1 and ri == 0) else ""

    # Актуальный адрес офиса на слайде контактов
    for slide in prs.slides:
        for s in slide.shapes:
            if not (hasattr(s, 'text') and s.text and 'Абая' in s.text):
                continue
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if 'Абая' in r.text:
                        r.text = COMPANY_ADDRESS


def fix_title_colors(prs):
    """Титульный лист: текст без явного цвета рендерится черным — делаем белым."""
    for s in prs.slides[0].shapes:
        if not (hasattr(s, 'text_frame') and s.has_text_frame):
            continue
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() and r.font.color.type is None:
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _delete_table_row(table, row_idx):
    tr = table.rows[row_idx]._tr
    tr.getparent().remove(tr)


def _clone_table_row(table, src_idx):
    """Клонирует строку таблицы (с форматированием) и добавляет в конец."""
    src_tr = table.rows[src_idx]._tr
    new_tr = copy.deepcopy(src_tr)
    src_tr.getparent().append(new_tr)
    return len(table.rows) - 1


def _estimate_text_lines(text: str) -> int:
    lines = 0
    for chunk in text.split('\n'):
        lines += max(1, math.ceil(len(chunk) / CASE_TABLE_CHARS_PER_LINE))
    return lines


# Эталонные позиции с кейс-слайдов 1-2 шаблона
STANDARD_CASE_TABLE_TOP = 2_444_005
STANDARD_CASE_TITLE_TOP = 1_937_774
STANDARD_CASE_PRIO_TOP = 1_955_844


def normalize_case_slide_layout(slide):
    """Выравнивает кейс-слайд по эталону: в шаблоне у слайда кейса 3
    заголовок и таблица сидят выше, чем у остальных."""
    table_shape = get_table_shape(slide)
    if table_shape is None:
        return
    delta = STANDARD_CASE_TABLE_TOP - table_shape.top
    if abs(delta) < 50_000:
        return
    for s in slide.shapes:
        # Блок "Ответ ИТ-руководства" внизу и картинки (их позиционируем отдельно) не трогаем
        if s.shape_type in (MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE):
            continue
        s.top = s.top + delta

    # Остаточная разница: заголовок/плашка приоритета на слайде кейса 3
    # имеют иной отступ от таблицы, чем на эталонных слайдах
    prio = get_priority(slide)
    if prio is not None:
        adj = STANDARD_CASE_PRIO_TOP - prio.top
        if adj:
            cx = prio.left + prio.width // 2
            cy = prio.top + prio.height // 2
            for s in slide.shapes:
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and \
                        s.left <= cx <= s.left + s.width and s.top <= cy <= s.top + s.height:
                    s.top = s.top + adj
            prio.top = STANDARD_CASE_PRIO_TOP

    title = get_title(slide)
    if title is not None and hasattr(title, 'top'):
        adj = STANDARD_CASE_TITLE_TOP - title.top
        if adj:
            for s in slide.shapes:
                if s.shape_type == MSO_SHAPE_TYPE.LINE:
                    s.top = s.top + adj
            title.top = STANDARD_CASE_TITLE_TOP


def clear_case_pictures(slide):
    """Удаляет шаблонные фото кейса (фотографии объектов прошлых клиентов).
    Возвращает центр X последней удаленной картинки для позиционирования новой."""
    old_center_x = None
    for s in list(slide.shapes):
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            old_center_x = s.left + s.width // 2
            _delete_shape(s)
    return old_center_x


def layout_case_image(slide, img_path):
    """Ставит картинку кейса под таблицей, ужимая так, чтобы не налезать
    на текст сверху и блок 'Ответ ИТ-руководства' снизу."""
    table_shape = get_table_shape(slide)

    # Нижняя граница зоны: группа "Ответ ИТ-руководства"
    bottom_limit = None
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            bottom_limit = s.top
            break
    if bottom_limit is None:
        bottom_limit = int(prs_height_fallback := 8_300_000)

    # Верхняя граница: оценка фактической высоты таблицы с текстом
    top_limit = 4_350_883  # позиция картинки в шаблоне как fallback
    old_center_x = clear_case_pictures(slide)

    if table_shape is not None:
        t = table_shape.table
        est = 0
        for r in range(len(t.rows)):
            est += _estimate_text_lines(t.cell(r, 0).text) * CASE_TABLE_LINE_EMU
            est += CASE_TABLE_ROW_PAD_EMU
        top_limit = table_shape.top + est

    available = bottom_limit - CASE_IMAGE_MARGIN_EMU - (top_limit + CASE_IMAGE_MARGIN_EMU)
    size = max(CASE_IMAGE_MIN_EMU, min(CASE_IMAGE_MAX_EMU, available))

    slide_width = 7_559_675
    center_x = old_center_x if old_center_x else slide_width // 2
    left = int(center_x - size // 2)
    top = int(top_limit + CASE_IMAGE_MARGIN_EMU)
    slide.shapes.add_picture(img_path, Emu(left), Emu(top), Emu(size), Emu(size))


def build_pptx(data: AuditData, audit_type: str = "full", auditor=None) -> io.BytesIO:
    prs = Presentation(TEMPLATE_PATH)
    include_recommendations = audit_type == "full"

    temp_images = []
    for i, case in enumerate(data.cases):
        if case.image_b64 and case.image_b64.startswith("data:image"):
            img_data = base64.b64decode(case.image_b64.split(",")[1])
            img_path = f"/tmp/case_img_{i}_{int(time.time())}.jpg"
            with open(img_path, "wb") as f:
                f.write(img_data)
            temp_images.append((i, img_path))

    try:
        remove_personal_marks(prs, auditor=auditor)
        fix_title_colors(prs)

        for s in prs.slides[0].shapes:
            if hasattr(s, 'text') and s.text:
                if 'Клиент' in s.text or 'Курылыс' in s.text:
                    perfect_replace(s, f"Клиент: {data.client_name}")

        for s in prs.slides[2].shapes:
            if hasattr(s, 'text') and s.text and 'аудита' in s.text:
                perfect_replace(s, data.review)

        prio_counts = {"ПЕРВЫЙ ПРИОРИТЕТ": 0, "ВТОРОЙ ПРИОРИТЕТ": 0, "ТРЕТИЙ ПРИОРИТЕТ": 0}

        # Категории в макете без нумерации — римские номера проставляем здесь,
        # по порядку появления. В таблице шаблона максимум 3 строки:
        # лишние категории считаем в последней.
        ordered_cats = []
        for case in data.cases:
            cat = _strip_category_number(case.category) or "Прочее"
            if cat not in ordered_cats:
                ordered_cats.append(cat)
        numbered = {name: f"{ROMAN[min(i, len(ROMAN)-1)]}. {name}" for i, name in enumerate(ordered_cats)}
        table_cats = [numbered[n] for n in ordered_cats[:3]]
        cat_counts = {c: {"ПЕРВЫЙ ПРИОРИТЕТ": 0, "ВТОРОЙ ПРИОРИТЕТ": 0, "ТРЕТИЙ ПРИОРИТЕТ": 0} for c in table_cats}

        for case in data.cases:
            prio_counts[case.priority] += 1
            cat = numbered[_strip_category_number(case.category) or "Прочее"]
            key = cat if cat in cat_counts else table_cats[-1]
            cat_counts[key][case.priority] += 1

        t3 = get_table(prs.slides[3])
        if t3:
            replace_table_cell(t3, 0, 0, str(prio_counts["ПЕРВЫЙ ПРИОРИТЕТ"]))
            replace_table_cell(t3, 0, 1, str(prio_counts["ВТОРОЙ ПРИОРИТЕТ"]))
            replace_table_cell(t3, 0, 2, str(prio_counts["ТРЕТИЙ ПРИОРИТЕТ"]))

        t4 = get_table(prs.slides[4])
        if t4:
            cats = table_cats
            for r, cat in enumerate(cats):
                replace_table_cell(t4, r, 0, cat)
                replace_table_cell(t4, r, 1, str(cat_counts[cat]["ПЕРВЫЙ ПРИОРИТЕТ"]))
                replace_table_cell(t4, r, 2, str(cat_counts[cat]["ВТОРОЙ ПРИОРИТЕТ"]))
                replace_table_cell(t4, r, 3, str(cat_counts[cat]["ТРЕТИЙ ПРИОРИТЕТ"]))

            # Лишние строки удаляем целиком, чтобы таблица не "плыла"
            for r in range(len(t4.rows) - 1, len(cats) - 1, -1):
                _delete_table_row(t4, r)

        conc_text = "\n".join([f"• {c}" for c in data.conclusions])
        for s in prs.slides[16].shapes:
            if hasattr(s, 'text') and s.text and 'Вариант 1' in s.text:
                perfect_replace(s, conc_text)

        # --- Секции и кейсы: клонируем шаблонные слайды под любое число кейсов ---
        divider_template = prs.slides[5]
        case_template = prs.slides[6]
        normalize_case_slide_layout(case_template)

        # группируем кейсы по категориям, сохраняя порядок появления категорий
        groups = {cat: [] for cat in ordered_cats}
        for i, case in enumerate(data.cases):
            cat = _strip_category_number(case.category) or "Прочее"
            groups[cat].append((i, case))

        case_no = 0
        for cat in ordered_cats:
            divider = clone_slide(prs, divider_template)
            for s in divider.shapes:
                if hasattr(s, 'text') and s.text and re.match(r"\s*[IVX]+\.", s.text):
                    perfect_replace(s, numbered[cat])

            for orig_idx, case in groups[cat]:
                case_no += 1
                slide = clone_slide(prs, case_template)
                perfect_replace(get_title(slide), f"Кейс {case_no}: {case.title}")

                if get_priority(slide):
                    perfect_replace(get_priority(slide), case.priority)
                    recolor_priority_badge(slide, case.priority)

                t = get_table(slide)
                if t:
                    replace_labeled_cell(t, 0, 0, "Уязвимость: ", case.vulnerability)
                    replace_labeled_cell(t, 1, 0, "Риски: ", case.risk)
                    if include_recommendations and case.recommendation:
                        rec_row = _clone_table_row(t, 1)
                        replace_labeled_cell(t, rec_row, 0, "Рекомендации: ", case.recommendation)

                img_path_tuple = next((x for x in temp_images if x[0] == orig_idx), None)
                if img_path_tuple:
                    layout_case_image(slide, img_path_tuple[1])
                else:
                    # Без новой картинки шаблонное фото чужого клиента в отчете не оставляем
                    clear_case_pictures(slide)

        # Убираем все шаблонные слайды кейсов/секций/предложений (5-15) —
        # вместо них уже добавлены клоны в конце
        for i in range(15, 4, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        # Порядок: [титул..категории] [клоны секций/кейсов] [выводы] [контакты]
        # После удаления выводы и контакты стоят перед клонами — переносим их в конец
        sld_lst = prs.slides._sldIdLst
        for _ in range(2):
            el = sld_lst[5]
            sld_lst.remove(el)
            sld_lst.append(el)

        output_io = io.BytesIO()
        prs.save(output_io)
        output_io.seek(0)
        return output_io
    finally:
        for _, img_path in temp_images:
            if os.path.exists(img_path):
                os.remove(img_path)
